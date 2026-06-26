# -*- coding: utf-8 -*-
"""Generate styled interbank rollover decks (EN + ZH): flow, rules, I/O only."""
from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

DOCS = Path(__file__).resolve().parents[1] / "docs"

SLATE_900 = RGBColor(0x0F, 0x17, 0x2A)
SLATE_700 = RGBColor(0x33, 0x41, 0x55)
SLATE_500 = RGBColor(0x64, 0x74, 0x8B)
SLATE_200 = RGBColor(0xE2, 0xE8, 0xF0)
SLATE_50 = RGBColor(0xF8, 0xFA, 0xFC)
SKY_500 = RGBColor(0x0E, 0xA5, 0xE9)
SKY_100 = RGBColor(0xE0, 0xF2, 0xFE)
EMERALD = RGBColor(0x10, 0xB9, 0x81)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


@dataclass
class LocaleContent:
    font: str
    footer: str
    cover_title: str
    cover_sub: str
    cover_meta: str
    cover_pill: str
    cover_tag: str
    timeline_title: str
    timeline_sub: str
    timeline_steps: list[tuple[str, str]]
    timeline_note: str
    coupon_title: str
    coupon_sub: str
    coupon_blocks: list[tuple[str, str]]
    coupon_note: str
    rules_title: str
    rules: list[str]
    io_title: str
    io_sub: str
    inputs_title: str
    inputs: list[str]
    outputs_title: str
    outputs: list[str]


LOCALES: dict[str, LocaleContent] = {
    "en": LocaleContent(
        font="Segoe UI",
        footer="Interbank Rollover",
        cover_title="Interbank\nRollover",
        cover_sub="Flow · Rules · Inputs & Outputs",
        cover_meta="Module: interbank_installment_rollover.py\nPer-step settlement in bank simulation",
        cover_pill="Quick reference",
        cover_tag="Process  →  coupon_cleared  →  Rules  →  I/O",
        timeline_title="Rollover × Step1–4 mapping",
        timeline_sub="Each column: source step (decentralized comments) + which rollover piece runs there",
        timeline_steps=[
            (
                "Step1 · Trade",
                "Source: #1) Trade\n_register_trade_contracts\n\nRollover here:\n• Report — read banks_in_rollover /\n  rollover_blocked (from last Step4)\n• Book — choose_trade_schedule\n  bullet vs installment\n• Register new contracts",
            ),
            (
                "Step2 · Intentions",
                "Source: #4) Step2\n collect_intentions\n\nRollover here:\n• Cap — rollover_borrow_quantity\n  (policy coupon_cleared):\n  coupon due → must clear\n  this step + proj cap>0;\n  no due → project cap only",
            ),
            (
                "Step3 · Match",
                "Source: #5) Step3\n RFQ / bipartite\n\nRollover here:\n• Match guard —\n  filter_borrowers_for_\n  rollover_block; RFQ &\n  trade落地 reject blocked\n  borrowers",
            ),
            (
                "Step4 · Settle",
                "Source: #7) Step4\n EN + recovery\n settle_interbank_period\n\nRollover here:\n• Settle — coupon pay;\n  bullet payoff\n• Renew — bullet →\n  installment if on\n• Report — coupon_due /\n  coupon_cleared sets;\n  rollover_blocked",
            ),
        ],
        timeline_note="Loop per period t: Step4ₜ → (blocked, log) → Step1ₜ₊₁ report/book → Step2ₜ₊₁ block → Step3ₜ₊₁ match → Step4ₜ₊₁ settle. (In simulate_step, Step4 may run before Step1–3 in the same t — same mapping.)",
        coupon_title="Borrow policy: coupon_cleared (default)",
        coupon_sub="rollover_borrow_quantity · checked each Step2 (intentions) / matching",
        coupon_blocks=[
            (
                "Who is in scope",
                "Banks with outstanding installment debt (rollover_blocked / "
                "active_installment_rollover_borrowers).",
            ),
            (
                "Not in rollover",
                "Normal borrow: liquidity need + project-opportunity need (unchanged).",
            ),
            (
                "Rollover — no coupon due this step",
                "May borrow only up to project-opportunity cap "
                "(compute_project_investment_borrow_cap). "
                "Liquidity-gap interbank borrow is blocked.",
            ),
            (
                "Rollover — coupon due this step",
                "May borrow only if: (1) borrower ∈ coupon_cleared "
                "(all installment coupons due this step paid via EN), "
                "and (2) project cap > 0. Otherwise quantity = 0.",
            ),
            (
                "Step4 → Step2 (same period t)",
                "settle_interbank_period fills coupon_due_borrowers & "
                "coupon_cleared_borrowers; simulate_step runs settlement "
                "before collect_intentions / bipartite match.",
            ),
        ],
        coupon_note="Constant: ROLLOVER_BORROW_COUPON_CLEARED · Function: rollover_borrow_quantity in interbank_installment_rollover.py",
        rules_title="Basic rules",
        rules=[
            "Contracts: bullet (overnight, one-shot payoff) or installment (20–120 business days, interest-first each period).",
            "New trade schedule (auto): large principal or low LCR → installment; small + high LCR → bullet.",
            "Rollover: when rollover_mode is on, bullet maturity repays then re-opens same principal as new installment (20–120 business days).",
            "Installment repayment: each coupon clears interest, then principal slice, via EN clearing on the lender–borrower pair.",
            "Borrow policy (default coupon_cleared): rollover borrowers with a coupon due this step may borrow only after that coupon clears via EN and project borrow cap > 0; no due → project cap only; liquidity-gap borrow blocked.",
            "Repayment funding: borrower liquid_assets + EN clearing only (optional central-bank liquidity backstop on shortfall).",
        ],
        io_title="Inputs & outputs",
        io_sub="Per call to settle_interbank_period and the matching phase",
        inputs_title="Inputs",
        inputs=[
            "step — current simulation period (int)",
            "ContractBook — active contracts (lender, borrower, schedule_type, principal, rates, periods_paid)",
            "banks[] — liquid_assets, liabilities, LCR-related fields, is_active",
            "ScheduleConfig — rollover_mode, tenor 20–120 business days, spreads, schedule_selection, principal thresholds",
            "New trades (matching) — amount, rate, borrower; blocked if borrower ∈ rollover_blocked_borrowers",
        ],
        outputs_title="Outputs",
        outputs=[
            "Step4: updated ContractBook; failed[]; log_rollover_status; rollover_blocked_borrowers",
            "Step1: new contracts; uses prior Step4 blocked set for reporting",
            "Step2–3: no new loans to blocked borrowers",
            "Bank state — interbank balances & cash after EN clearing (Step4)",
        ],
    ),
    "zh": LocaleContent(
        font="Microsoft YaHei",
        footer="同业 Rollover",
        cover_title="同业\nRollover",
        cover_sub="流程 · 规则 · 输入与输出",
        cover_meta="模块：interbank_installment_rollover.py\n银行仿真每步同业 rollover 流程",
        cover_pill="速查",
        cover_tag="流程  →  coupon_cleared  →  规则  →  输入/输出",
        timeline_title="Rollover 与 Step1–4 对应",
        timeline_sub="每列：源码步骤（decentralized 注释）+ 该步执行的 rollover 环节",
        timeline_steps=[
            (
                "Step1 · 成交",
                "源码：#1) Trade\n_register_trade_contracts\n\nRollover：\n• 统计 — 读取 banks_in_rollover /\n  rollover_blocked（上一步 Step4 产出）\n• 选型 — choose_trade_schedule\n  bullet / installment\n• 入账 — 新合约写入簿",
            ),
            (
                "Step2 · 意图",
                "源码：#4) Step2\n collect_intentions\n\nRollover：\n• 额度 — rollover_borrow_quantity\n  （默认 coupon_cleared）：\n  本期有到期 → 须还清\n  且项目融资能力>0；\n  无到期 → 仅项目额度",
            ),
            (
                "Step3 · 撮合",
                "源码：#5) Step3\n RFQ / 双边撮合\n\nRollover：\n• 撮合拦截 —\n  filter_borrowers；\n  RFQ 与成交登记\n  拒绝 blocked 借款人",
            ),
            (
                "Step4 · 结算",
                "源码：#7) Step4\n EN + recovery\n settle_interbank_period\n\nRollover：\n• 结算 — 分期先息后本；\n  bullet 到期结清\n• 续借 — bullet→installment\n• 统计 — coupon_due /\n  coupon_cleared；\n  rollover_blocked",
            ),
        ],
        timeline_note="周期 t：Step4ₜ 结算并统计 → Step1ₜ₊₁ 读状态并入账 → Step2ₜ₊₁ 封锁 → Step3ₜ₊₁ 撮合 → Step4ₜ₊₁ 结算。（simulate_step 同期 t 内可能先执行 Step4 再 Step1–3，对应关系不变。）",
        coupon_title="借款政策：coupon_cleared（默认）",
        coupon_sub="rollover_borrow_quantity · 每步 Step2 意图 / 撮合时检查",
        coupon_blocks=[
            (
                "适用对象",
                "仍有未还清同业分期的借款人（rollover_blocked / "
                "active_installment_rollover_borrowers）。",
            ),
            (
                "非 rollover 银行",
                "按原逻辑：流动性缺口需求 + 项目机会需求均可计入借款意图。",
            ),
            (
                "rollover · 本期无分期到期",
                "仅允许借至项目融资额度上限（compute_project_investment_borrow_cap）；"
                "禁止用同业新借补流动性缺口。",
            ),
            (
                "rollover · 本期有分期到期",
                "仅当：(1) 借款人 ∈ coupon_cleared（本期所有到期合约经 EN 付清），"
                "且 (2) 项目融资额度 > 0，才允许按项目额度借；否则额度为 0。",
            ),
            (
                "Step4 → Step2（同期 t）",
                "settle_interbank_period 产出 coupon_due_borrowers、"
                "coupon_cleared_borrowers；simulate_step 先结算再生成意图/撮合。",
            ),
        ],
        coupon_note="常量 ROLLOVER_BORROW_COUPON_CLEARED · 函数 rollover_borrow_quantity（interbank_installment_rollover.py）",
        rules_title="基本规则",
        rules=[
            "合约类型：bullet（隔夜、到期一次结清）或 installment（20–120 个工作日，每期先息后本）。",
            "成交选型（auto）：本金大或 LCR 低 → installment；本金小且 LCR 高 → bullet。",
            "Rollover：rollover_mode 开启时，bullet 到期先结清，再以相同本金续为 20–120 个工作日 installment。",
            "分期偿付：每期先清算利息、再清算本金份额，对借贷双方走 EN 清算。",
            "借款政策（默认 coupon_cleared）：本期有分期到期的 rollover 借款人，须当期息/本经 EN 还清且项目融资额度>0 才可借；本期无到期仅允许项目额度；禁止用同业新借补流动性缺口。",
            "还款资金：借款人 liquid_assets + EN 清算（流动性不足时可触发央行 backstop）。",
        ],
        io_title="输入与输出",
        io_sub="settle_interbank_period 与撮合阶段的接口",
        inputs_title="输入 Input",
        inputs=[
            "step — 当前仿真期数",
            "ContractBook — 存续合约（对手方、schedule_type、本金、利率、已付期数等）",
            "banks[] — 流动性、负债、LCR 相关字段、is_active",
            "ScheduleConfig — rollover_mode、期限 20–120 个工作日、利差、选型阈值等",
            "新成交（撮合）— 金额、利率、借款人；borrower ∈ rollover_blocked 则拒绝",
        ],
        outputs_title="输出 Output",
        outputs=[
            "Step4：更新 ContractBook；failed[]；log_rollover_status；rollover_blocked_borrowers",
            "Step1：新合约；读取上一步 Step4 的 blocked 集合用于统计",
            "Step2–3：blocked 银行无新增借款",
            "银行状态 — Step4 清算后同业科目与现金",
        ],
    ),
}

OUTPUTS = {
    "en": [DOCS / "interbank_rollover_overview_styled.pptx", DOCS / "interbank_rollover_overview.pptx"],
    "zh": [DOCS / "interbank_rollover_overview_styled_zh.pptx", DOCS / "interbank_rollover_overview_zh.pptx"],
}


class Deck:
    def __init__(self, loc: LocaleContent):
        self.loc = loc
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.n = 0
        self._saved_path: Path | None = None

    def save(self, paths: list[Path]) -> Path:
        DOCS.mkdir(parents=True, exist_ok=True)
        for path in paths:
            try:
                self.prs.save(str(path))
                self._saved_path = path
                return path
            except PermissionError:
                continue
        raise PermissionError(f"Cannot write PPT — close {paths[0].name} in PowerPoint and rerun.")

    def _font(self) -> str:
        return self.loc.font

    def _slide(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def _rect(self, slide, l, t, w, h, fill, line=None, radius=None):
        kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        sh = slide.shapes.add_shape(kind, l, t, w, h)
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
        if line:
            sh.line.color.rgb = line
            sh.line.width = Pt(1)
        else:
            sh.line.fill.background()
        return sh

    def _text(self, slide, l, t, w, h, text, size=18, bold=False, color=SLATE_700, align=PP_ALIGN.LEFT):
        box = slide.shapes.add_textbox(l, t, w, h)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = self._font()
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = align
        return box

    def _sidebar(self, slide):
        self._rect(slide, Inches(0), Inches(0), Inches(0.22), self.prs.slide_height, SKY_500)

    def _footer(self, slide):
        self.n += 1
        self._rect(slide, Inches(0.22), Inches(7.05), Inches(13.0), Inches(0.01), SLATE_200)
        self._text(slide, Inches(0.55), Inches(7.12), Inches(8), Inches(0.3), self.loc.footer, size=9, color=SLATE_500)
        self._text(slide, Inches(12.35), Inches(7.08), Inches(0.7), Inches(0.35), str(self.n), size=11, bold=True, color=SKY_500, align=PP_ALIGN.RIGHT)

    def _title_bar(self, slide, title: str, subtitle: str | None = None):
        self._rect(slide, Inches(0.22), Inches(0), Inches(13.11), Inches(1.15), SLATE_900)
        self._text(slide, Inches(0.55), Inches(0.22), Inches(10), Inches(0.55), title, size=30, bold=True, color=WHITE)
        if subtitle:
            self._text(slide, Inches(0.55), Inches(0.72), Inches(10), Inches(0.35), subtitle, size=13, color=SLATE_200)

    def slide_cover(self):
        loc = self.loc
        s = self._slide()
        self._rect(s, Inches(0), Inches(0), Inches(7.2), self.prs.slide_height, SLATE_900)
        for x, y, r, alpha in [(8.8, 1.2, 1.8, SKY_100), (10.5, 4.5, 2.2, SKY_100), (7.5, 5.8, 1.0, SLATE_200)]:
            c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(r), Inches(r))
            c.fill.solid()
            c.fill.fore_color.rgb = alpha
            c.line.fill.background()
        self._text(s, Inches(0.65), Inches(1.6), Inches(6.2), Inches(1.5), loc.cover_title, size=48, bold=True, color=WHITE)
        self._text(s, Inches(0.65), Inches(3.35), Inches(6.0), Inches(0.9), loc.cover_sub, size=20, color=SLATE_200)
        self._text(s, Inches(0.65), Inches(4.4), Inches(5.8), Inches(1.0), loc.cover_meta, size=13, color=SLATE_500)
        pill = self._rect(s, Inches(0.65), Inches(5.85), Inches(2.0), Inches(0.42), SKY_500, radius=True)
        pill.text_frame.paragraphs[0].text = loc.cover_pill
        pill.text_frame.paragraphs[0].font.name = self._font()
        pill.text_frame.paragraphs[0].font.size = Pt(11)
        pill.text_frame.paragraphs[0].font.bold = True
        pill.text_frame.paragraphs[0].font.color.rgb = WHITE
        pill.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        self._text(s, Inches(7.6), Inches(2.8), Inches(5.2), Inches(2.0), loc.cover_tag, size=22, bold=True, color=SLATE_700)
        self.n += 1

    def slide_timeline(self):
        loc = self.loc
        s = self._slide()
        self._sidebar(s)
        self._title_bar(s, loc.timeline_title, loc.timeline_sub)
        self._footer(s)
        y_line = Inches(3.35)
        self._rect(s, Inches(0.9), y_line, Inches(11.5), Inches(0.06), SLATE_200)
        span = 11.0 / (len(loc.timeline_steps) - 1)
        card_w = 2.35 if len(loc.timeline_steps) <= 4 else 1.9
        for i, (name, desc) in enumerate(loc.timeline_steps):
            cx = 0.9 + i * span
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - 0.18), Inches(3.18), Inches(0.42), Inches(0.42))
            dot.fill.solid()
            dot.fill.fore_color.rgb = SKY_500
            dot.line.fill.background()
            dot.text_frame.paragraphs[0].text = str(i + 1)
            dot.text_frame.paragraphs[0].font.name = self._font()
            dot.text_frame.paragraphs[0].font.size = Pt(11)
            dot.text_frame.paragraphs[0].font.bold = True
            dot.text_frame.paragraphs[0].font.color.rgb = WHITE
            dot.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            half = card_w / 2
            self._rect(s, Inches(cx - half), Inches(4.0), Inches(card_w), Inches(2.35), WHITE, SLATE_200, radius=True)
            tf = s.shapes.add_textbox(Inches(cx - half + 0.08), Inches(4.08), Inches(card_w - 0.16), Inches(2.2)).text_frame
            tf.paragraphs[0].text = name
            tf.paragraphs[0].font.name = self._font()
            tf.paragraphs[0].font.size = Pt(12)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = SLATE_900
            p2 = tf.add_paragraph()
            p2.text = desc
            p2.font.name = self._font()
            p2.font.size = Pt(9)
            p2.font.color.rgb = SLATE_500
        self._rect(s, Inches(0.55), Inches(6.55), Inches(12.2), Inches(0.75), SKY_100, radius=True)
        self._text(s, Inches(0.75), Inches(6.68), Inches(11.8), Inches(0.55), loc.timeline_note, size=10, color=SLATE_700)

    def slide_rules(self):
        loc = self.loc
        s = self._slide()
        self._sidebar(s)
        self._title_bar(s, loc.rules_title)
        self._footer(s)
        y = Inches(1.45)
        for i, rule in enumerate(loc.rules):
            accent = SKY_500 if i % 2 == 0 else EMERALD
            self._rect(s, Inches(0.55), y, Inches(0.08), Inches(0.85), accent)
            self._rect(s, Inches(0.7), y, Inches(12.05), Inches(0.85), WHITE, SLATE_200, radius=True)
            self._text(s, Inches(0.9), y + Inches(0.1), Inches(11.7), Inches(0.7), rule, size=13, color=SLATE_700)
            y += Inches(0.95)

    def slide_coupon_cleared(self):
        loc = self.loc
        s = self._slide()
        self._sidebar(s)
        self._title_bar(s, loc.coupon_title, loc.coupon_sub)
        self._footer(s)
        y = Inches(1.42)
        row_h = Inches(0.98)
        for i, (heading, body) in enumerate(loc.coupon_blocks):
            accent = SKY_500 if i % 2 == 0 else EMERALD
            self._rect(s, Inches(0.55), y, Inches(0.08), row_h - Inches(0.08), accent)
            self._rect(s, Inches(0.7), y, Inches(12.05), row_h - Inches(0.08), WHITE, SLATE_200, radius=True)
            self._text(
                s, Inches(0.9), y + Inches(0.06), Inches(11.7), Inches(0.28),
                heading, size=13, bold=True, color=SLATE_900,
            )
            self._text(
                s, Inches(0.9), y + Inches(0.34), Inches(11.7), Inches(0.58),
                body, size=11, color=SLATE_700,
            )
            y += row_h
        self._rect(s, Inches(0.55), Inches(6.52), Inches(12.2), Inches(0.78), SKY_100, radius=True)
        self._text(s, Inches(0.75), Inches(6.65), Inches(11.8), Inches(0.55), loc.coupon_note, size=10, color=SLATE_700)

    def slide_io(self):
        loc = self.loc
        s = self._slide()
        self._sidebar(s)
        self._title_bar(s, loc.io_title, loc.io_sub)
        self._footer(s)
        self._rect(s, Inches(6.55), Inches(1.45), Inches(0.02), Inches(5.35), SLATE_200)

        def col(x, title, items, accent):
            self._rect(s, x, Inches(1.45), Inches(5.9), Inches(0.5), accent, radius=True)
            box = s.shapes[-1]
            box.text_frame.paragraphs[0].text = title
            box.text_frame.paragraphs[0].font.name = self._font()
            box.text_frame.paragraphs[0].font.size = Pt(15)
            box.text_frame.paragraphs[0].font.bold = True
            box.text_frame.paragraphs[0].font.color.rgb = WHITE
            box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            y = Inches(2.1)
            for line in items:
                self._rect(s, x + Inches(0.05), y, Inches(5.8), Inches(0.78), WHITE, SLATE_200, radius=True)
                self._text(s, x + Inches(0.2), y + Inches(0.1), Inches(5.5), Inches(0.62), f"• {line}", size=12, color=SLATE_700)
                y += Inches(0.86)

        col(Inches(0.55), loc.inputs_title, loc.inputs, SKY_500)
        col(Inches(6.85), loc.outputs_title, loc.outputs, SLATE_700)


def build_deck(locale: str) -> Deck:
    d = Deck(LOCALES[locale])
    d.slide_cover()
    d.slide_timeline()
    d.slide_coupon_cleared()
    d.slide_rules()
    d.slide_io()
    return d


def main():
    saved = []
    for locale in ("en", "zh"):
        path = build_deck(locale).save(OUTPUTS[locale])
        saved.append(f"{locale}: {path.name}")
    print("Saved PPT in docs/ (5 slides: cover + flow + coupon_cleared + rules + I/O)")
    for line in saved:
        print(" ", line)


if __name__ == "__main__":
    main()
